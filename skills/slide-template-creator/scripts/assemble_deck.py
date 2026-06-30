"""assemble_deck.py — Bolt the template's brand bookends onto a content deck.

The content slides are built by another skill (pptxgenjs / html2pptx), saved as
a standalone .pptx. This script takes that deck, prepends `temp01.pptx`'s brand
cover slide (with your title/subtitle) and appends its closing / "thank you"
slide, and writes a single merged file — preserving the template's master,
logos and colour bands on the bookends.

Usage:
    python scripts/assemble_deck.py CONTENT.pptx \
        --title "簡報標題" [--subtitle "副標"] [--lang zh-tw|en] \
        [--template temp01.pptx] [--out OUTPUT.pptx]

How it works:
    base = clean temp01  →  cover (native, brand)  →  imported content slides
    →  closing (native, brand)  →  save.

Content slides are cloned element-by-element with their image/media
relationships remapped, and geometry + font sizes uniformly scaled if the
content deck's slide size differs from the template (both assumed 16:9).
"""
import argparse
import copy
import os
import sys

from pptx import Presentation
from pptx.oxml.ns import qn

# Import the brand helpers from the skill root.
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from slide_helpers import (  # noqa: E402
    load_clean_template,
    find_layout_by_role,
    add_cover_slide,
    add_closing_slide,
)

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ── Blank content canvas ─────────────────────────────────────────────────────
def find_blank_layout(prs):
    """Pick a placeholder-free layout to host imported content slides.

    Prefer a layout literally named 'Blank'; otherwise the one with the fewest
    placeholders. Imported slides carry all their own visuals, so we want the
    emptiest possible canvas (master background still shows through, which is
    fine for full-bleed content)."""
    named = [l for l in prs.slide_layouts if "blank" in l.name.lower()]
    pool = named or list(prs.slide_layouts)
    pool.sort(key=lambda l: len(list(l.placeholders)))
    chosen = pool[0]
    print(f"[LAYOUT] content host → '{chosen.name}' "
          f"({len(list(chosen.placeholders))} placeholders)")
    return chosen


# ── Relationship remapping ───────────────────────────────────────────────────
_REL_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"), qn("r:pict"),
              qn("r:dm"), qn("r:lo"), qn("r:qs"), qn("r:cs"))


def _remap_rels(el, src_part, dest_part):
    """Rewrite every relationship id in `el` (and descendants) so it points at
    a copy of the source part registered in `dest_part`."""
    cache = {}
    for node in el.iter():
        for attr in _REL_ATTRS:
            rid = node.get(attr)
            if not rid:
                continue
            if rid in cache:
                node.set(attr, cache[rid])
                continue
            try:
                rel = src_part.rels[rid]
            except KeyError:
                continue
            if rel.is_external:
                new_rid = dest_part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True)
            else:
                new_rid = dest_part.relate_to(rel.target_part, rel.reltype)
            cache[rid] = new_rid
            node.set(attr, new_rid)


# ── Uniform scaling (handles slide-size mismatch) ────────────────────────────
def _scale_tree(el, f):
    """Scale geometry (a:off/a:ext/a:chOff/a:chExt), line widths (a:ln@w) and
    font sizes (@sz) by factor f, in place."""
    if abs(f - 1.0) < 1e-6:
        return
    for node in el.iter():
        tag = node.tag
        if tag in (qn("a:off"), qn("a:chOff")):
            for k in ("x", "y"):
                v = node.get(k)
                if v is not None:
                    node.set(k, str(int(round(int(v) * f))))
        elif tag in (qn("a:ext"), qn("a:chExt")):
            for k in ("cx", "cy"):
                v = node.get(k)
                if v is not None:
                    node.set(k, str(int(round(int(v) * f))))
        elif tag == qn("a:ln"):
            w = node.get("w")
            if w is not None:
                node.set("w", str(int(round(int(w) * f))))
        # font sizes live on any rPr/defRPr/endParaRPr as @sz
        sz = node.get("sz")
        if sz is not None and tag in (
                qn("a:rPr"), qn("a:defRPr"), qn("a:endParaRPr")):
            node.set("sz", str(max(100, int(round(int(sz) * f)))))


# ── Slide cloning ────────────────────────────────────────────────────────────
def clone_slide(src_slide, src_prs, dest_prs, dest_layout, scale):
    """Append a copy of `src_slide` to `dest_prs` on `dest_layout`."""
    dest_slide = dest_prs.slides.add_slide(dest_layout)
    dest_tree = dest_slide.shapes._spTree

    # Drop any shapes/placeholders the layout injected — start from empty canvas.
    for shp in list(dest_slide.shapes):
        shp._element.getparent().remove(shp._element)

    for shp in src_slide.shapes:
        new_el = copy.deepcopy(shp._element)
        _remap_rels(new_el, src_slide.part, dest_slide.part)
        _scale_tree(new_el, scale)
        dest_tree.append(new_el)

    # Carry over an explicit slide background, if the source set one.
    csld = src_slide._element.find(qn("p:cSld"))
    if csld is not None:
        bg = csld.find(qn("p:bg"))
        if bg is not None:
            new_bg = copy.deepcopy(bg)
            _remap_rels(new_bg, src_slide.part, dest_slide.part)
            dest_csld = dest_slide._element.find(qn("p:cSld"))
            dest_csld.insert(0, new_bg)
    return dest_slide


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("content", help="content deck (.pptx) from the content skill")
    ap.add_argument("--title", required=True, help="cover title")
    ap.add_argument("--subtitle", default="", help="cover subtitle")
    ap.add_argument("--lang", default="en", help="en | zh-tw")
    ap.add_argument("--template", default=None,
                    help="brand template (default: temp01.pptx next to the skill)")
    ap.add_argument("--out", default=None, help="output path (default: <content>_final.pptx)")
    args = ap.parse_args()

    skill_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    template = args.template or os.path.join(skill_root, "temp01.pptx")
    out = args.out or os.path.splitext(args.content)[0] + "_final.pptx"

    print("[ASSEMBLE] slide-template-creator — bookend mode")
    print(f"  template : {template}")
    print(f"  content  : {args.content}")
    print(f"  output   : {out}")

    content_prs = Presentation(args.content)
    prs = load_clean_template(template)

    # Scale factor: match content geometry to the template canvas (16:9 assumed).
    scale = prs.slide_width / content_prs.slide_width
    if abs(scale - 1.0) > 1e-6:
        print(f"  [scale] content {content_prs.slide_width.inches:.3f}\" → "
              f"template {prs.slide_width.inches:.3f}\"  (×{scale:.4f})")

    # 1. Brand cover.
    add_cover_slide(prs, args.title, args.subtitle, lang=args.lang)

    # 2. Imported content slides.
    blank = find_blank_layout(prs)
    n = 0
    for s in content_prs.slides:
        clone_slide(s, content_prs, prs, blank, scale)
        n += 1
    print(f"  [import] {n} content slide(s) cloned")

    # 3. Brand closing / thank-you.
    add_closing_slide(prs)

    prs.save(out)
    print(f"[ASSEMBLE] saved {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
