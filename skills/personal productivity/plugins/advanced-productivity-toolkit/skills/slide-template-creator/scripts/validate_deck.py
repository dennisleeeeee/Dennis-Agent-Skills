"""validate_deck.py — Lightweight post-build deck audit.

Usage:
    python scripts/validate_deck.py <path-to-output.pptx>

Heuristic structural checks. The agent should treat output as a guide, not a
gate — a slide can land its takeaway in many ways. Hard failures are reserved
for things that are unambiguously broken (empty title, single-shape slide).
Exit code is non-zero only if at least one HARD failure fires; soft warnings
print but don't fail the build.

Rules:
  HARD  shape_count < 2                        → slide is empty
  HARD  title text length == 0 on a content slide
  SOFT  no callout / box in lower 30% of slide → takeaway placement check
  SOFT  any body placeholder > 60 chars        → looks like bullet fallback
  SOFT  title text length <= 12 chars          → likely topic label, not insight
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def audit(path: str) -> int:
    prs = Presentation(path)
    slide_h = prs.slide_height
    hard_failures = 0
    soft_warnings = 0

    for i, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_text_lens = []
        shape_count = 0
        shape_in_lower_third = False

        for shape in slide.shapes:
            shape_count += 1
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title_text = txt
                elif shape.is_placeholder:
                    body_text_lens.append(len(txt))
            if shape.top is not None and shape.top > slide_h * 0.70:
                shape_in_lower_third = True

        is_content_slide = shape_count > 3
        if not is_content_slide:
            print(f"slide {i:02d}  [skip — brand layer (cover/section/closing)]")
            continue

        hard = []
        soft = []
        if shape_count < 2:
            hard.append("empty slide (shape_count < 2)")
        if len(title_text) == 0:
            hard.append("empty title on content slide")
        if not shape_in_lower_third:
            soft.append("no element in lower 30% — check takeaway placement")
        if any(n > 60 for n in body_text_lens):
            soft.append("placeholder body text > 60 chars (bullet fallback?)")
        if 0 < len(title_text) <= 12:
            soft.append(f"title short ('{title_text}') — likely topic label, not insight")

        if hard:
            hard_failures += 1
            print(f"slide {i:02d}  ❌ HARD: " + "; ".join(hard)
                  + ("  | ⚠  " + "; ".join(soft) if soft else ""))
        elif soft:
            soft_warnings += 1
            print(f"slide {i:02d}  ⚠  " + "; ".join(soft))
        else:
            print(f"slide {i:02d}  ✅ ok")

    print()
    if hard_failures:
        print(f"{hard_failures} HARD failure(s), {soft_warnings} warning(s).")
        return 1
    if soft_warnings:
        print(f"All slides pass HARD checks. {soft_warnings} soft warning(s) — review and decide.")
    else:
        print("All slides pass.")
    return 0


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("usage: python scripts/validate_deck.py <deck.pptx>")
    sys.exit(audit(sys.argv[1]))
