"""slide_helpers.py — Reusable visual primitives for the corporate template.

Import in any slide-generation script:
    from slide_helpers import *

This module loads the corporate template (temp01.pptx), strips its demo slides,
and provides 10 visual patterns + low-level primitives for drawing branded slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette ────────────────────────────────────────────────────────────
BLUE     = RGBColor(0x00, 0x78, 0xD4)
NAVY     = RGBColor(0x24, 0x3A, 0x5E)
CYAN     = RGBColor(0x50, 0xE6, 0xFF)
TEAL     = RGBColor(0x00, 0x85, 0x75)
DARKTEAL = RGBColor(0x27, 0x4B, 0x47)
GRAY     = RGBColor(0x73, 0x73, 0x73)
LIGHT    = RGBColor(0xE6, 0xE6, 0xE6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x00, 0x00, 0x00)
SOFT_BG  = RGBColor(0xF3, 0xF6, 0xFA)
RED      = RGBColor(0xD8, 0x3B, 0x01)

SLIDE_W = 13.333
SLIDE_H = 7.5

# ── Template inspection (run BEFORE load_clean_template) ─────────────────────
def inspect_template(path):
    """
    Analyze a .pptx template and print its key properties.
    Call this FIRST — before load_clean_template() — to discover:
      • Available layout names (pick the right ones for each slide role)
      • Slide dimensions
      • Number of demo slides to be stripped

    Returns a dict. Use the 'layouts' list to choose layout names when calling
    add_slide() or find_layout_by_role().

    Usage:
        info = inspect_template("my_template.pptx")
        prs  = load_clean_template("my_template.pptx")
        slide = prs.slides.add_slide(find_layout_by_role(prs, "content", info))
    """
    prs = Presentation(path)
    layouts = [l.name for l in prs.slide_layouts]
    w = round(prs.slide_width.inches, 3)
    h = round(prs.slide_height.inches, 3)
    demo_count = len(prs.slides)

    # Try to read theme accent colors from slide master XML
    theme_colors = {}
    try:
        master = prs.slide_masters[0]
        theme_el = master.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}theme')
        if theme_el is None:
            # Try alternate namespace path
            for el in master.element.iter():
                if el.tag.endswith('}theme'):
                    theme_el = el
                    break
        if theme_el is not None:
            accent_names = ['dk1','lt1','dk2','lt2','accent1','accent2','accent3',
                            'accent4','accent5','accent6','hlink','folHlink']
            for name in accent_names:
                for srgb in theme_el.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
                    parent_tag = srgb.getparent().tag.split('}')[-1] if srgb.getparent() is not None else ''
                    if parent_tag == name:
                        theme_colors[name] = '#' + srgb.get('val','')
    except Exception:
        pass

    print(f"[TEMPLATE INSPECT] ─────────────────────────────")
    print(f"  File       : {path}")
    print(f"  Slide size : {w}\" × {h}\"  ({round(w*914400)} × {round(h*914400)} EMU)")
    print(f"  Demo slides: {demo_count}  (will be stripped by load_clean_template)")
    print(f"  Layouts ({len(layouts)}):")
    for i, name in enumerate(layouts):
        print(f"    [{i:3d}] {name}")
    if theme_colors:
        print(f"  Theme colors: {theme_colors}")
    print(f"[TEMPLATE INSPECT] ─────────────────────────────")

    return {
        "path": path,
        "slide_w": w,
        "slide_h": h,
        "demo_count": demo_count,
        "layouts": layouts,
        "theme_colors": theme_colors,
    }


def find_layout_by_role(prs, role, template_info=None):
    """
    Find the best layout for a semantic role, without hardcoding layout names.
    Role values: 'cover' | 'section' | 'content' | 'closing'

    Strategy:
      1. Match by common name patterns for that role
      2. If no match, fall back to placeholder-count heuristic:
         cover → 2 placeholders, content → 1 placeholder (title only), closing → 0 or 1
      3. Last resort: first layout

    Always print which layout was selected so the caller can verify.
    """
    role_patterns = {
        'cover':   ['title slide', 'cover', 'title and subtitle'],
        'section': ['section title', 'section header', 'section', 'divider'],
        'content': ['title only', 'blank', 'title, content', 'title and content'],
        'closing': ['closing', 'end slide', 'thank you', 'end', 'last'],
    }
    patterns = role_patterns.get(role, [])

    # 1. Name-based match (case-insensitive)
    for pattern in patterns:
        for layout in prs.slide_layouts:
            if pattern.lower() in layout.name.lower():
                print(f"[LAYOUT] role='{role}' → '{layout.name}' (name match)")
                return layout

    # 2. Heuristic: count placeholders
    if role == 'cover':
        # Title + subtitle = 2 placeholders
        for layout in prs.slide_layouts:
            phs = list(layout.placeholders)
            if len(phs) >= 2:
                print(f"[LAYOUT] role='{role}' → '{layout.name}' (2-ph heuristic)")
                return layout
    elif role == 'content':
        # Title Only = 1 placeholder (idx 0)
        for layout in prs.slide_layouts:
            phs = list(layout.placeholders)
            if len(phs) == 1 and phs[0].placeholder_format.idx == 0:
                print(f"[LAYOUT] role='{role}' → '{layout.name}' (title-only heuristic)")
                return layout

    # 3. Last resort
    fallback = prs.slide_layouts[0]
    print(f"[LAYOUT] role='{role}' → '{fallback.name}' (fallback)")
    return fallback


# ── Template & layout management ─────────────────────────────────────────────
def load_clean_template(path):
    prs = Presentation(path)
    xs = prs.slides._sldIdLst
    while len(xs) > 0:
        rId = xs[0].get(qn('r:id'))
        del xs[0]
        prs.part.drop_rel(rId)
    return prs

def get_layout(prs, name):
    for l in prs.slide_layouts:
        if l.name == name:
            return l
    return prs.slide_layouts[0]

def add_slide(prs, layout_name):
    return prs.slides.add_slide(get_layout(prs, layout_name))

def set_layout_title(slide, text, keep_master_font=False):
    """Fill the title placeholder.

    keep_master_font=True → don't override the font (use this for cover, section
    divider, closing — anywhere brand fidelity matters or text is Chinese).
    keep_master_font=False (default) → force Segoe UI Semibold, suitable for
    Latin-only content slides.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            if not keep_master_font:
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = "Segoe UI Semibold"
            return

def hide_body_phs(slide, keep_idx=None):
    """Remove non-title placeholders. Pass keep_idx=(12,) to preserve a
    subtitle/body placeholder you actually want to fill (e.g. cover subtitle).
    Default behavior unchanged: removes everything except idx 0.
    """
    keep = set(keep_idx or ())
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx == 0 or idx in keep:
            continue
        ph._element.getparent().remove(ph._element)

# ── Brand-fidelity helpers (cover / section / closing) ───────────────────────
# These use the master template's own styling instead of overriding fonts.
# Use them for the brand layer; use the pattern helpers below for content.

# Chinese-capable fallback chain. python-pptx writes one font name per run, so
# pick the one that renders on the target machine. On macOS/Windows both are
# bundled; Linux servers should embed PingFang or Noto Sans TC.
CN_TITLE_FONT = "Microsoft JhengHei UI"
CN_BODY_FONT  = "Microsoft JhengHei"

def add_cover_slide(prs, title, subtitle="", lang="en"):
    """Create the brand cover slide. Fills title + subtitle placeholders and
    preserves master styling (logo, background, color band). Do NOT decorate.

    lang='zh-tw' → use Microsoft JhengHei UI; lang='en' → keep master font.
    """
    layout = find_layout_by_role(prs, "cover")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
            if lang.startswith("zh"):
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = CN_TITLE_FONT
        elif subtitle and (idx == 12 or ph.placeholder_format.type == 2):
            ph.text = subtitle
            if lang.startswith("zh"):
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = CN_BODY_FONT
    return slide

def add_section_slide(prs, title, lang="en"):
    """Brand section divider. Title-only placeholder, master decoration."""
    layout = find_layout_by_role(prs, "section")
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            if lang.startswith("zh"):
                for p in ph.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = CN_TITLE_FONT
            break
    return slide

def add_closing_slide(prs):
    """Brand closing slide. Logo/decoration only — no text. Caller should
    add a small text box if a closing message is wanted, but the master is
    designed to stand alone.
    """
    layout = find_layout_by_role(prs, "closing")
    return prs.slides.add_slide(layout)

# ── Text helpers ─────────────────────────────────────────────────────────────
def set_text(tf, text, font="Segoe UI", size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04);  tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = align; p.text = text
    for run in p.runs:
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color

def add_text_box(slide, left, top, w, h, text, **kw):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    set_text(tb.text_frame, text, **kw)
    return tb

def write_bullets(slide, x, y, w, h, lines, size=12, color=BLACK, bullet=True, gap=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0);  tf.margin_bottom = Inches(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = (f"\u2022  {line}" if bullet and line else line)
        p.space_after = Pt(gap)
        for r in p.runs:
            r.font.name = "Segoe UI"; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

# ── Shape primitives ─────────────────────────────────────────────────────────
def add_card(slide, left, top, w, h, fill=WHITE, line=LIGHT, line_w=0.75, rounded=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.5):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width)
    ln = line.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle'); tailEnd.set('w', 'med'); tailEnd.set('len', 'med')
    return line

# ── Pattern 1: Stat Cards ────────────────────────────────────────────────────
def draw_stat_cards(slide, stats, y=1.55, h=1.7, gap=0.3):
    n = len(stats)
    total_w = SLIDE_W - 2.0
    w = (total_w - gap*(n-1)) / n
    x0 = (SLIDE_W - (w*n + gap*(n-1))) / 2
    for i, s in enumerate(stats):
        x = x0 + i*(w+gap)
        add_card(slide, x, y, w, h, fill=WHITE, line=LIGHT)
        tb = slide.shapes.add_textbox(Inches(x+0.1), Inches(y+0.05),
                                        Inches(w*0.65), Inches(1.1))
        set_text(tb.text_frame, s["num"], font="Segoe UI Semibold",
                  size=64, bold=True, color=s["color"], align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
        tb = slide.shapes.add_textbox(Inches(x+w*0.65+0.05), Inches(y+0.55),
                                        Inches(w*0.3), Inches(0.5))
        set_text(tb.text_frame, s["unit"], size=18, bold=True, color=s["color"])
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+1.25),
                                        Inches(w-0.4), Inches(0.4))
        set_text(tb.text_frame, s["desc"], size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ── Pattern 2: Flow Pipeline ─────────────────────────────────────────────────
def draw_flow_pipeline(slide, items, colors, y=4.05, h=0.75, gap=0.18):
    n = len(items)
    total_w = SLIDE_W - 2.0
    w = (total_w - gap*(n-1)) / n
    x0 = (SLIDE_W - (w*n + gap*(n-1))) / 2
    for i, text in enumerate(items):
        x = x0 + i*(w+gap)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = colors[i]
        s.line.fill.background()
        set_text(s.text_frame, text, size=13, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n-1:
            add_arrow(slide, x+w, y+h/2, x+w+gap, y+h/2, color=GRAY, width=1.5)

# ── Pattern 3: Numbered Content Cards ────────────────────────────────────────
def draw_numbered_cards(slide, cards, y=1.55, h=4.2, gap=0.3):
    n = len(cards)
    total_w = SLIDE_W - 2.0
    w = (total_w - gap*(n-1)) / n
    x0 = (SLIDE_W - (w*n + gap*(n-1))) / 2
    for i, c in enumerate(cards):
        x = x0 + i*(w+gap)
        add_card(slide, x, y, w, h, fill=WHITE, line=LIGHT)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x), Inches(y),
                                        Inches(w), Inches(0.95))
        band.fill.solid(); band.fill.fore_color.rgb = c["color"]
        band.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.05),
                                        Inches(1.1), Inches(0.9))
        set_text(tb.text_frame, c["num"], font="Segoe UI Semibold",
                  size=42, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        tb = slide.shapes.add_textbox(Inches(x+1.35), Inches(y+0.12),
                                        Inches(w-1.5), Inches(0.4))
        set_text(tb.text_frame, c["title"], font="Segoe UI Semibold",
                  size=16, bold=True, color=WHITE)
        tb = slide.shapes.add_textbox(Inches(x+1.35), Inches(y+0.5),
                                        Inches(w-1.5), Inches(0.35))
        set_text(tb.text_frame, c["sub"], size=12, color=WHITE)
        write_bullets(slide, x+0.3, y+1.15, w-0.6, h-1.3,
                       c["body"], size=13, color=BLACK, gap=8)

# ── Pattern 4: Comparison Table ──────────────────────────────────────────────
def draw_comparison_table(slide, left_header, right_header, rows,
                           y=1.55, header_h=0.65, row_h=0.78, col_gap=0.5):
    cw = (SLIDE_W - 2.0 - col_gap) / 2
    lx = (SLIDE_W - (cw*2 + col_gap)) / 2
    rx = lx + cw + col_gap
    hl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(lx), Inches(y), Inches(cw), Inches(header_h))
    hl.fill.solid(); hl.fill.fore_color.rgb = GRAY; hl.line.fill.background()
    set_text(hl.text_frame, left_header, size=16, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    hr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(rx), Inches(y), Inches(cw), Inches(header_h))
    hr.fill.solid(); hr.fill.fore_color.rgb = BLUE; hr.line.fill.background()
    set_text(hr.text_frame, right_header, size=16, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ry = y + header_h
    for i, (label, old, new) in enumerate(rows):
        bg = SOFT_BG if i % 2 == 0 else WHITE
        for x_pos, txt, txt_color in [(lx, old, BLACK), (rx, new, NAVY)]:
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(x_pos), Inches(ry+i*row_h),
                                            Inches(cw), Inches(row_h))
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.line.color.rgb = LIGHT; cell.line.width = Pt(0.5)
            tb = slide.shapes.add_textbox(
                Inches(x_pos+0.2), Inches(ry+i*row_h+0.07),
                Inches(cw-0.4), Inches(0.3))
            set_text(tb.text_frame, label, size=10, bold=True, color=GRAY)
            tb = slide.shapes.add_textbox(
                Inches(x_pos+0.2), Inches(ry+i*row_h+0.34),
                Inches(cw-0.4), Inches(row_h-0.35))
            set_text(tb.text_frame, txt, size=14, bold=True, color=txt_color)
    arr_y = y + header_h + (len(rows)*row_h)/2 - 0.3
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(lx+cw+0.05), Inches(arr_y),
                                   Inches(col_gap-0.1), Inches(0.6))
    arr.fill.solid(); arr.fill.fore_color.rgb = NAVY
    arr.line.fill.background()

# ── Pattern 5: Hub-and-Spoke ─────────────────────────────────────────────────
def draw_hub_spoke(slide, hub_text, nodes, hub_color=NAVY, hub_size=1.8,
                    cx=None, cy=3.95):
    cx = cx or SLIDE_W/2
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx-hub_size/2), Inches(cy-hub_size/2),
                                   Inches(hub_size), Inches(hub_size))
    hub.fill.solid(); hub.fill.fore_color.rgb = hub_color
    hub.line.fill.background()
    set_text(hub.text_frame, hub_text, font="Segoe UI Semibold",
              size=22, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cw, ch = 3.5, 1.5
    for n in nodes:
        x = cx + n["dx"] - cw/2
        y = cy + n["dy"] - ch/2
        x1 = x + cw if n["dx"] < 0 else x
        x2 = cx + (-hub_size/2 if n["dx"] < 0 else hub_size/2)
        add_arrow(slide, x1, y+ch/2, x2, cy, color=LIGHT, width=2)
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(cw), Inches(ch))
        c.fill.solid(); c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = n["color"]; c.line.width = Pt(2)
        p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x+0.2), Inches(y+0.15),
                                    Inches(0.75), Inches(0.4))
        p.adjustments[0] = 0.5
        p.fill.solid(); p.fill.fore_color.rgb = n["color"]
        p.line.fill.background()
        set_text(p.text_frame, n["label"], size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb = slide.shapes.add_textbox(Inches(x+1.05), Inches(y+0.18),
                                        Inches(cw-1.2), Inches(0.35))
        set_text(tb.text_frame, n["sub"], size=11, color=GRAY)
        tb = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.7),
                                        Inches(cw-0.4), Inches(0.7))
        set_text(tb.text_frame, n["title"], font="Segoe UI Semibold",
                  size=16, bold=True, color=NAVY)

# ── Pattern 6: Mock UI Side-by-Side ──────────────────────────────────────────
def draw_mock_ui_comparison(slide, left, right, y=1.55, h=4.5, gap=0.5):
    pw = (SLIDE_W - 2.0 - gap) / 2
    plx = (SLIDE_W - (pw*2 + gap)) / 2
    prx = plx + pw + gap

    def _draw_panel(x, panel):
        is_dash = panel["type"] == "dashboard"
        add_card(slide, x, y, pw, h, fill=WHITE,
                  line=panel["header_color"] if is_dash else LIGHT,
                  line_w=1.5 if is_dash else 0.75)
        hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x), Inches(y), Inches(pw), Inches(0.55))
        hb.fill.solid(); hb.fill.fore_color.rgb = panel["header_color"]
        hb.line.fill.background()
        set_text(hb.text_frame, panel["title"], font="Segoe UI Semibold",
                  size=14, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if panel["type"] == "text_wall":
            my = y + 0.8
            for i in range(13):
                line_len = pw - 0.6 - (0.6 if i%4==0 else 0) - (1.4 if i in [5,9] else 0)
                ml = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(x+0.3), Inches(my+i*0.22),
                                              Inches(line_len), Inches(0.1))
                ml.fill.solid(); ml.fill.fore_color.rgb = LIGHT
                ml.line.fill.background()
        else:
            ux = x + 0.3; uy = y + 0.8
            mock_stats = panel.get("stats", [("42%","Conv.",CYAN),("8.7s","Load",BLUE),("1.2k","Users",NAVY)])
            for i, (val, lbl, color) in enumerate(mock_stats):
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               Inches(ux+i*1.7), Inches(uy),
                                               Inches(1.55), Inches(0.9))
                box.fill.solid(); box.fill.fore_color.rgb = color
                box.line.fill.background()
                tb = slide.shapes.add_textbox(Inches(ux+i*1.7+0.1), Inches(uy+0.05),
                                                Inches(1.35), Inches(0.55))
                set_text(tb.text_frame, val, font="Segoe UI Semibold",
                          size=24, bold=True, color=WHITE,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                tb = slide.shapes.add_textbox(Inches(ux+i*1.7+0.1), Inches(uy+0.55),
                                                Inches(1.35), Inches(0.3))
                set_text(tb.text_frame, lbl, size=10, color=WHITE, align=PP_ALIGN.CENTER)
            chy = uy + 1.05
            cbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(ux), Inches(chy),
                                           Inches(pw-0.6), Inches(1.6))
            cbg.fill.solid(); cbg.fill.fore_color.rgb = SOFT_BG
            cbg.line.color.rgb = LIGHT
            bar_h = [0.6, 1.0, 0.8, 1.3, 1.1, 0.9, 1.4]
            bw = (pw-0.6-0.4) / len(bar_h)
            for i, bh in enumerate(bar_h):
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               Inches(ux+0.2+i*bw),
                                               Inches(chy+1.5-bh),
                                               Inches(bw*0.7), Inches(bh))
                bar.fill.solid()
                bar.fill.fore_color.rgb = BLUE if i%2==0 else CYAN
                bar.line.fill.background()
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x+0.3), Inches(y+h-0.85),
                                       Inches(pw-0.6), Inches(0.65))
        lbl.fill.solid(); lbl.fill.fore_color.rgb = SOFT_BG
        lbl.line.color.rgb = panel["label_color"]; lbl.line.width = Pt(1)
        set_text(lbl.text_frame, panel["label"], size=12, color=panel["label_color"],
                  bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _draw_panel(plx, left)
    _draw_panel(prx, right)

# ── Pattern 7: Icon Concept Cards ────────────────────────────────────────────
def draw_icon_cards(slide, cards, y=1.55, h=4.2, gap=0.3):
    n = len(cards)
    total_w = SLIDE_W - 2.0
    w = (total_w - gap*(n-1)) / n
    x0 = (SLIDE_W - (w*n + gap*(n-1))) / 2
    for i, c in enumerate(cards):
        x = x0 + i*(w+gap)
        add_card(slide, x, y, w, h, fill=WHITE, line=LIGHT)
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y), Inches(w), Inches(1.2))
        top.fill.solid(); top.fill.fore_color.rgb = c["color"]
        top.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+0.1),
                                        Inches(1.0), Inches(1.0))
        set_text(tb.text_frame, c["icon"], size=44,
                  anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        tb = slide.shapes.add_textbox(Inches(x+1.3), Inches(y+0.2),
                                        Inches(w-1.4), Inches(0.4))
        set_text(tb.text_frame, c["tag"], font="Segoe UI Semibold",
                  size=20, bold=True, color=WHITE)
        tb = slide.shapes.add_textbox(Inches(x+1.3), Inches(y+0.62),
                                        Inches(w-1.4), Inches(0.35))
        set_text(tb.text_frame, c["title"], size=11, color=WHITE)
        tb = slide.shapes.add_textbox(Inches(x+0.25), Inches(y+1.3),
                                        Inches(w-0.4), Inches(0.4))
        set_text(tb.text_frame, c["desc"], font="Segoe UI Semibold",
                  size=13, bold=True, color=c["color"])
        write_bullets(slide, x+0.25, y+1.8, w-0.4, h-2.0,
                       c["ex"], size=12, color=BLACK, bullet=False, gap=6)

# ── Pattern 8: Quote Slide ───────────────────────────────────────────────────
def draw_quote_slide(slide, quote, author, source):
    qm = slide.shapes.add_textbox(Inches(1.5), Inches(1.4),
                                    Inches(2.5), Inches(2.5))
    set_text(qm.text_frame, '\u201C', size=220, bold=True, color=BLUE)
    qt = slide.shapes.add_textbox(Inches(2.5), Inches(2.5),
                                    Inches(9.5), Inches(2.5))
    set_text(qt.text_frame, quote, font="Segoe UI Semibold",
              size=28, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(2.5), Inches(5.5),
                                   Inches(0.08), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(2.8), Inches(5.5), Inches(9.0), Inches(1.0))
    set_text(tb.text_frame, author, font="Segoe UI Semibold",
              size=18, bold=True, color=NAVY)
    tb = slide.shapes.add_textbox(Inches(2.8), Inches(5.95),
                                    Inches(9.0), Inches(0.5))
    set_text(tb.text_frame, source, size=14, color=GRAY)

# ── Pattern 9: Stop / Start Cards ────────────────────────────────────────────
def draw_stop_start_rows(slide, items, y=1.6, row_h=1.35, gap=0.2):
    for i, (stop, start, num) in enumerate(items):
        ry = y + i*(row_h+gap)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(0.8), Inches(ry+row_h/2-0.45),
                                         Inches(0.9), Inches(0.9))
        badge.fill.solid(); badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        set_text(badge.text_frame, num, font="Segoe UI Semibold",
                  size=20, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        sx = 2.0; sw = 4.8
        sc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(sx), Inches(ry), Inches(sw), Inches(row_h))
        sc.fill.solid(); sc.fill.fore_color.rgb = SOFT_BG
        sc.line.color.rgb = GRAY; sc.line.width = Pt(0.75)
        tb = slide.shapes.add_textbox(Inches(sx+0.2), Inches(ry+0.15),
                                        Inches(sw-0.4), Inches(0.4))
        set_text(tb.text_frame, "\u2717  STOP", font="Segoe UI Semibold",
                  size=11, bold=True, color=GRAY)
        tb = slide.shapes.add_textbox(Inches(sx+0.2), Inches(ry+0.5),
                                        Inches(sw-0.4), Inches(row_h-0.55))
        set_text(tb.text_frame, stop, size=14, color=BLACK, anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(slide, sx+sw+0.05, ry+row_h/2, sx+sw+0.55, ry+row_h/2,
                   color=BLUE, width=2.5)
        stx = sx + sw + 0.7
        stc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(stx), Inches(ry),
                                       Inches(sw), Inches(row_h))
        stc.fill.solid(); stc.fill.fore_color.rgb = WHITE
        stc.line.color.rgb = BLUE; stc.line.width = Pt(2)
        tb = slide.shapes.add_textbox(Inches(stx+0.2), Inches(ry+0.15),
                                        Inches(sw-0.4), Inches(0.4))
        set_text(tb.text_frame, "\u2713  START", font="Segoe UI Semibold",
                  size=11, bold=True, color=BLUE)
        tb = slide.shapes.add_textbox(Inches(stx+0.2), Inches(ry+0.5),
                                        Inches(sw-0.4), Inches(row_h-0.55))
        set_text(tb.text_frame, start, size=14, bold=True, color=NAVY,
                  anchor=MSO_ANCHOR.MIDDLE)

# ── Pattern 10: Callout Bar ──────────────────────────────────────────────────
def draw_callout_bar(slide, text, y=6.45, h=0.7, style="solid_navy"):
    if style == "source_bar":
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(1.0), Inches(y),
                                      Inches(SLIDE_W-2.0), Inches(h))
        bg.fill.solid(); bg.fill.fore_color.rgb = SOFT_BG
        bg.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1.0), Inches(y),
                                       Inches(0.1), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.3), Inches(y),
                                        Inches(SLIDE_W-2.5), Inches(h))
        set_text(tb.text_frame, text, size=12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    elif style == "soft_blue":
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.0), Inches(y),
                                     Inches(SLIDE_W-2.0), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = SOFT_BG
        s.line.color.rgb = BLUE; s.line.width = Pt(1.5)
        set_text(s.text_frame, text, size=14, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.0), Inches(y),
                                     Inches(SLIDE_W-2.0), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = NAVY
        s.line.fill.background()
        set_text(s.text_frame, text, font="Segoe UI Semibold",
                  size=14, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── Pattern 11: Text + Right Visual Panel (split layout) ─────────────────────
def draw_text_visual_split(slide, text_blocks, visual_items, y=1.55, h=4.55,
                            split=0.47, text_x=0.8, gap=0.35):
    """
    Left 47% = rich text content; right panel = visual evidence (stat cards or
    bullet highlights). Mirrors the side-by-side style of the AI reference deck.

    text_blocks : list of dicts  {"label": str, "body": str, "color": RGBColor}
                  label = small bold category label above body text
    visual_items: list of dicts  {"num": str, "unit": str, "desc": str, "color": RGBColor}
                  (same format as draw_stat_cards rows — displayed vertically)
    """
    text_w = (SLIDE_W - text_x - gap) * split
    vis_x  = text_x + text_w + gap
    vis_w  = SLIDE_W - vis_x - 0.5

    # ── left: stacked text blocks ──────────────────────────────────────────
    block_h = h / max(len(text_blocks), 1)
    for i, blk in enumerate(text_blocks):
        by = y + i * block_h
        # accent bar on left edge
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(text_x), Inches(by + 0.12),
                                      Inches(0.06), Inches(block_h - 0.28))
        bar.fill.solid(); bar.fill.fore_color.rgb = blk.get("color", BLUE)
        bar.line.fill.background()
        # label
        tb = slide.shapes.add_textbox(Inches(text_x + 0.18), Inches(by + 0.12),
                                       Inches(text_w - 0.22), Inches(0.28))
        set_text(tb.text_frame, blk.get("label", ""), size=10, bold=True,
                  color=blk.get("color", BLUE))
        # body
        tb = slide.shapes.add_textbox(Inches(text_x + 0.18), Inches(by + 0.38),
                                       Inches(text_w - 0.22), Inches(block_h - 0.48))
        set_text(tb.text_frame, blk.get("body", ""), size=13, color=BLACK)
        tb.text_frame.word_wrap = True

    # ── right: vertical stat stack ────────────────────────────────────────
    item_h = h / max(len(visual_items), 1)
    for i, item in enumerate(visual_items):
        iy = y + i * item_h
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(vis_x), Inches(iy + 0.1),
                                       Inches(vis_w), Inches(item_h - 0.2))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = item.get("color", BLUE)
        card.line.width = Pt(2)
        card.shadow.inherit = False
        # colored top band
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(vis_x), Inches(iy + 0.1),
                                       Inches(vis_w), Inches(0.28))
        band.fill.solid(); band.fill.fore_color.rgb = item.get("color", BLUE)
        band.line.fill.background()
        # number
        num_w = vis_w * 0.55
        tb = slide.shapes.add_textbox(Inches(vis_x + 0.12), Inches(iy + 0.36),
                                       Inches(num_w), Inches(item_h * 0.45))
        set_text(tb.text_frame, item.get("num", ""), font="Segoe UI Semibold",
                  size=36, bold=True, color=item.get("color", BLUE),
                  anchor=MSO_ANCHOR.MIDDLE)
        # unit
        tb = slide.shapes.add_textbox(Inches(vis_x + num_w + 0.15), Inches(iy + 0.45),
                                       Inches(vis_w - num_w - 0.2), Inches(0.35))
        set_text(tb.text_frame, item.get("unit", ""), size=14, bold=True,
                  color=item.get("color", BLUE))
        # description
        tb = slide.shapes.add_textbox(Inches(vis_x + 0.12), Inches(iy + item_h * 0.57),
                                       Inches(vis_w - 0.24), Inches(item_h * 0.38))
        set_text(tb.text_frame, item.get("desc", ""), size=11, color=GRAY)


# ── Pattern 12: Section Label Header ─────────────────────────────────────────
def draw_section_label(slide, part_text, topic_text, x=0.8, y=1.1):
    """
    Draws a two-line section identifier above the slide title area, like:
        PART 1 · FOUNDATION
        AI 如何重塑企業工作流程

    part_text  e.g. "PART 1 · FOUNDATION"
    topic_text e.g. "AI 如何重塑企業工作流程"

    The slide title placeholder should still be set separately via set_layout_title().
    This label goes in the body area, above the main content block.
    """
    # small spaced label
    tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                   Inches(SLIDE_W - x - 0.5), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = part_text.upper()
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.name = "Segoe UI"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = GRAY
    # spacing between letters via tracking can't be done in pptx; use spaced string

    # accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y + 0.32),
                                   Inches(2.4), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # topic text
    tb2 = slide.shapes.add_textbox(Inches(x), Inches(y + 0.4),
                                    Inches(SLIDE_W - x - 0.5), Inches(0.65))
    set_text(tb2.text_frame, topic_text, font="Segoe UI Semibold",
              size=22, bold=True, color=NAVY)
